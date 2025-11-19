from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Load the template
prs = Presentation('Corporate-PPT-Template-2025.pptx')

# Function to add title slide (using Cover layout)
def add_title_slide(title, subtitle, date):
    slide_layout = prs.slide_layouts[0]  # Cover layout
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    placeholders = slide.placeholders
    if len(placeholders) > 1:
        placeholders[1].text = subtitle
    if len(placeholders) > 2:
        placeholders[2].text = date
    return slide

# Function to add agenda slide
def add_agenda_slide(title, items):
    slide_layout = prs.slide_layouts[1]  # Agenda layout
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        for item in items:
            p = tf.add_paragraph()
            p.text = item
    return slide

# Function to add content slide with table
def add_table_slide(title, headers, rows):
    slide_layout = prs.slide_layouts[2]  # Tables layout
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    # Add table (assume placeholder or add new)
    table = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(1), Inches(2), Inches(8), Inches(4)).table
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    return slide

# Function to add diagram slide (e.g., ZIA)
def add_diagram_slide(title):
    slide_layout = prs.slide_layouts[3]  # Diagrams layout
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    # Add boxes and connectors (example)
    shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(2), Inches(1))
    shape1.text = "User authentication and provisioning"
    # Add more shapes, connectors, etc., as per template
    return slide

# Add slides with content from Pixartprinting
add_title_slide("Professional Services Transition Meeting", "PIXARTPRINTING", "19/09/2025")
add_agenda_slide("Meeting Agenda", ["Project Summary", "Technical Summary", "Recommended Next Steps"])
add_title_slide("Project Summary")
# Add project status report slide with table
add_table_slide("Milestones", ["Milestone", "Baseline Date", "Target Completion Date", "Status"], milestone_defaults)  # Use defaults or dynamic
# Add more tables for rollout, objectives, deliverables, open items
add_diagram_slide("Deployed ZIA Architecture")
# Add next steps, thank you, etc.

# Save the new PPT
prs.save('generated_transition.pptx')
print("PPT generated successfully.")
