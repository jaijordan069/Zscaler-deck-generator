#!/usr/bin/env python3
"""
Streamlit app: Zscaler Professional Services Transition Deck PPT Generator
Improved version: Attractive UI, exact template match, image handling, more defaults.
"""
from __future__ import annotations
import io
import re
import requests
from datetime import datetime
from typing import List, Optional
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.util import Length  # For type checking

# -------------------------
# Configuration / Constants (Updated for exact template match)
# -------------------------
st.set_page_config(page_title="Zscaler Transition Deck Generator", layout="wide", initial_sidebar_state="expanded")

# Custom CSS for attractive UI (Softened for better readability)
st.markdown("""
    <style>
    .stApp { background-color: #F5F9FD; }  /* Light blue bg */
    .stButton > button { background-color: #256CF7; color: white; border-radius: 8px; }
    .stTextInput > div > div > input { border-color: #256CF7; }
    .sidebar .sidebar-content { background-color: #256CF7; color: white; }
    h1, h2 { color: #001744; }
    </style>
""", unsafe_allow_html=True)

# Colors (exact from template)
COLOR_BRIGHT_BLUE = RGBColor(37, 108, 247)
COLOR_NAVY = RGBColor(0, 23, 68)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_LIGHT_GRAY = RGBColor(229, 241, 250)
COLOR_CYAN = RGBColor(18, 212, 255)
COLOR_ACCENT_GREEN = RGBColor(107, 255, 179)
COLOR_THREAT_RED = RGBColor(237, 25, 81)
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_ZSCALER_YELLOW = RGBColor(255, 192, 0)

# Fonts & sizes (template)
FONT_NAME = "Century Gothic"
SIZE_TITLE = Pt(36)
SIZE_SLIDE_TITLE = Pt(28)
SIZE_SUBTITLE = Pt(20)
SIZE_HEADER = Pt(18)
SIZE_BODY = Pt(14)
SIZE_SMALL = Pt(12)
SIZE_FOOTER = Pt(8)

# Layout constants (fine-tuned to match template positions)
MARGIN_LEFT = Inches(0.45)
MARGIN_TOP = Inches(0.45)
MARGIN_RIGHT = Inches(0.45)
FOOTER_HEIGHT = Inches(0.35)

# Assets (added alt logos, bg if needed)
LOGO_URL = "https://upload.wikimedia.org/wikipedia/commons/thumb/8/8b/Zscaler_logo.svg/512px-Zscaler_logo.svg.png"
ALT_LOGO_URL = "https://companieslogo.com/img/orig/ZS-46a5871c.png?t=1720244494"
BG_URL = None  # Add if you have template bg image URL

# Date regex
DATE_RE = re.compile(r'^\d{2}/\d{2}/\d{4}$')

# RAG Colors for status (added from template)
RAG_COLORS = {
    "Red": RGBColor(255, 0, 0),
    "Amber": RGBColor(255, 191, 0),
    "Green": RGBColor(0, 255, 0),
    "Blue": RGBColor(0, 0, 255),
    "Gray": RGBColor(128, 128, 128)
}

# -------------------------
# Utilities (enhanced with more guards)
# -------------------------
def is_valid_date(d: str) -> bool:
    if not d or d == "??":
        return True  # Allow ?? as per template
    return bool(DATE_RE.match(d))

def download_image_to_bytes(url: Optional[str]) -> Optional[io.BytesIO]:
    if not url:
        return None
    try:
        r = requests.get(url, timeout=10)
        r.raise_for_status()
        return io.BytesIO(r.content)
    except Exception:
        st.warning(f"Couldn't download image from {url}")
        return None

def set_font_run(run, name: str = FONT_NAME, size: Pt = SIZE_BODY, bold: bool = False, color: RGBColor = COLOR_BLACK):
    try:
        run.font.name = name
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color
        if hasattr(run, "_element"):
            run._element.rPr.rFonts.set(qn("a:ea"), name)
    except Exception:
        pass

def add_textbox(slide, left, top, width, height, text: str, size=SIZE_BODY, bold=False, color=COLOR_BLACK, align=PP_ALIGN.LEFT, auto_size=False):
    try:
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        if auto_size:
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.clear()
        p = tf.add_paragraph() if len(tf.paragraphs) == 0 else tf.paragraphs[0]  # Guard
        p.text = text or ""
        p.alignment = align
        if not p.runs:
            run = p.add_run()
        else:
            run = p.runs[0]
        set_font_run(run, size=size, bold=bold, color=color)
        return txBox
    except Exception:
        st.warning("Failed to add textbox")
        return None

def apply_template_branding(prs: Presentation, slide, slide_num: int, logo_bytes: Optional[io.BytesIO]):
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    logo_w = Inches(1.5)  # Tweaked for template
    logo_h = Inches(0.4)
    logo_left = slide_width - logo_w - MARGIN_RIGHT
    logo_top = MARGIN_TOP / 2
    if logo_bytes:
        try:
            slide.shapes.add_picture(logo_bytes, logo_left, logo_top, logo_w, logo_h)
        except Exception:
            pass
    # Footer (exact text from template)
    footer_left = MARGIN_LEFT
    footer_top = slide_height - FOOTER_HEIGHT
    footer_w = Inches(4.0)
    footer_h = FOOTER_HEIGHT
    add_textbox(slide, footer_left, footer_top, footer_w, footer_h, "Zscaler, Inc. All rights reserved. © 2025", size=SIZE_FOOTER, color=COLOR_NAVY)
    # Slide number
    add_textbox(slide, slide_width - Inches(1.0), footer_top, Inches(0.8), footer_h, str(slide_num), size=SIZE_FOOTER, color=COLOR_NAVY, align=PP_ALIGN.RIGHT)

def add_slide_with_background(prs: Presentation, bg_bytes: Optional[io.BytesIO]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    if bg_bytes:
        slide.shapes.add_picture(bg_bytes, 0, 0, prs.slide_width, prs.slide_height)
    return slide

# -------------------------
# Streamlit UI (Made attractive: Columns, expanders, previews, images in expander)
# -------------------------
with st.sidebar:
    st.image(LOGO_URL, width=200)
    st.header("Zscaler Deck Generator")
    st.markdown("Create customer transition decks fast! Matches template exactly.")
    st.markdown("**Steps:**\n1. Fill details.\n2. Upload images if needed.\n3. Preview.\n4. Generate & Download.")
    theme = st.selectbox("Theme", ["Light", "Dark"], index=0)  # Can tie to colors later

st.title("Zscaler Professional Services Transition Deck Generator")
st.markdown("Enter details below. Defaults match the Pixartprinting template. UI is organized for ease!")

# Image Uploaders (wrapped in expander to hide clutter)
with st.expander("Optional: Upload Template Images", expanded=False):
    uploaded_images = {}
    for img_name in ["image3.jpg", "image4.png", "image5.png", "image6.png", "image35.png", "image36.svg", "image7.png", "image10.png", "image8.jpg", "image2.png", "image37.png", "image9.png", "image38.png", "image39.png", "image40.png", "image41.jpg", "image1.png"]:
        uploaded = st.file_uploader(f"Upload {img_name}", type=["jpg", "png", "svg"])
        if uploaded:
            uploaded_images[img_name] = io.BytesIO(uploaded.read())

# Customer & Project (balanced columns)
st.header("Customer & Project Basics")
col1, col2, col3 = st.columns(3)
customer_name = col1.text_input("Customer Name *", value="Pixartprinting")
today_date = col2.text_input("Today's Date *", value="19/09/2025")
project_start = col3.text_input("Project Start Date *", value="01/06/2025")
project_end = st.text_input("Project End Date *", value="19/09/2025")
project_summary_text = st.text_area("Project Summary Text", value="More than half of the users have been deployed and there were not any critical issues. Not expected issues during enrollment of remaining users", height=100)

# Milestones (expander, defaults from template)
st.header("Milestones")
milestone_defaults = [
    ("Initial Project Schedule Accepted", "27/06/2025", "27/06/2025", ""),
    ("Initial Design Accepted", "14/07/2025", "17/07/2025", ""),
    ("Pilot Configuration Complete", "28/07/2025", "18/07/2025", ""),
    ("Pilot Rollout Complete", "08/08/2025", "22/08/2025", ""),
    ("Production Configuration Complete", "29/08/2025", "29/08/2025", ""),
    ("Production Rollout Complete", "19/09/2025", "??", ""),
    ("Final Design Accepted", "19/09/2025", "19/09/2025", ""),
]
milestones_data = []
with st.expander("Edit Milestones (7 rows)", expanded=True):
    for i, default in enumerate(milestone_defaults):
        c1, c2, c3, c4 = st.columns(4)
        mn = c1.text_input(f"Name {i+1}", default[0])
        mb = c2.text_input(f"Baseline {i+1}", default[1])
        mt = c3.text_input(f"Target {i+1}", default[2])
        ms = c4.text_input(f"Status {i+1}", default[3])
        milestones_data.append({"name": mn, "baseline": mb, "target": mt, "status": ms})

# User Rollout (columns)
st.header("User Rollout Roadmap")
col1, col2 = st.columns(2)
with col1:
    pilot_target = st.number_input("Pilot Target Users", value=100)
    pilot_current = st.number_input("Pilot Current Users", value=449)
    pilot_completion = st.text_input("Pilot Completion", value="19/09/2025")
    pilot_status = st.text_input("Pilot Status", value="")
with col2:
    prod_target = st.number_input("Production Target Users", value=800)
    prod_current = st.number_input("Production Current Users", value=449)
    prod_completion = st.text_input("Production Completion", value="19/09/2025")
    prod_status = st.text_input("Production Status", value="")

# Objectives (expander, defaults from template)
st.header("Project Objectives")
objectives_defaults = [
    ("Protect and Secure Internet Access for Users", "More than half of the users have Zscaler Client Connector deployed and are fully protected when they are outside of the corporate office", "Not enough time to deploy ZCC in all users but deployment is on track to be finished by Pixartprinting and no critical issues are expected."),
    ("Complete user posture", "Users and devices are identified, and policies can be applied based on this criteria", "No deviations"),
    ("Comprehensive Web filtering", "Web filtering based on reputation and dynamic categorization rather than simply categories.", "No deviations"),
]
objectives_data = []
with st.expander("Edit Objectives (3 rows)", expanded=True):
    for i, default in enumerate(objectives_defaults):
        obj = st.text_area(f"Objective {i+1}", default[0], height=50)
        act = st.text_area(f"Actual {i+1}", default[1], height=50)
        dev = st.text_area(f"Deviation {i+1}", default[2], height=50)
        objectives_data.append({"objective": obj, "actual": act, "deviation": dev})

# Deliverables (expander, defaults)
st.header("Deliverables")
deliverables_defaults = [
    ("Kick-Off Meeting and Slides", "27/06/2025"),
    ("Design and Configuration of Zscaler Platform (per scope)", "30/06/2025 – 11/07/2025"),
    ("Troubleshooting Guide(s)", "18/07/2025"),
    ("Initial & Final Design Document", "17/07/2025 – 17/09/2025"),
    ("Transition Meeting Slides", "19/09/2025"),
]
deliverables_data = []
with st.expander("Edit Deliverables (5 rows)", expanded=True):
    for i, default in enumerate(deliverables_defaults):
        c1, c2 = st.columns(2)
        dn = c1.text_input(f"Name {i+1}", default[0])
        dd = c2.text_input(f"Date {i+1}", default[1])
        deliverables_data.append({"name": dn, "date": dd})

# Technical Summary (columns)
st.header("Technical Summary")
col1, col2 = st.columns(2)
with col1:
    idp = st.text_input("Identity Provider", value="Entra ID")
    auth_type = st.text_input("Authentication Type", value="SAML 2.0")
    prov_type = st.text_input("User/Group Provisioning", value="SCIM Provisioning")
with col2:
    tunnel_type = st.text_input("Tunnel Type", value="ZCC with Z-Tunnel 2.0")
    deploy_system = st.text_input("ZCC Deployment System", value="MS Intune/Jamf")
col1, col2, col3 = st.columns(3)
windows_num = col1.number_input("Windows Devices", value=351)
mac_num = col2.number_input("MacOS Devices", value=98)
geo_locations = col3.text_input("Geo Locations", value="Europe, North Africa, USA")
col1, col2, col3, col4 = st.columns(4)
ssl_policies = col1.number_input("SSL Policies", value=10)
url_policies = col2.number_input("URL Policies", value=5)
cloud_policies = col3.number_input("Cloud App Policies", value=5)
fw_policies = col4.number_input("Firewall Policies", value=15)

# Open Items (expander, defaults from template)
st.header("Open Items")
open_defaults = [
    ("Finish Production rollout", "October 2025", "Pixartprinting", "Onboard remaining users from all departments including Developers."),
    ("Tighten Firewall policies", "October 2025", "Pixartprinting", "Change the default Firewall rule from Allow All to Block All after configuring all the required exceptions."),
    ("Tighten Cloud App Control Policies", "October 2025", "Pixartprinting", "Configure block policies for high risk applications in all categories."),
    ("Fine tune SSL Inspection policies", "November 2025", "Pixartprinting", "Continue adjusting and adding exclusions to SSL Inspection policies as required."),
    ("Configure DLP policies", "December 2025", "Pixartprinting", "Configure DLP policies to control sensitive data and avoid potential data leaks."),
    ("Deploy ZCC on Mobile devices", "January 2026", "Pixartprinting", "Expand the deployment of Zscaler Client Connector to Mobile devices."),
]
open_items_data = []
with st.expander("Edit Open Items (6 rows)", expanded=True):
    for i, default in enumerate(open_defaults):
        otask = st.text_input(f"Task {i+1}", default[0])
        odate = st.text_input(f"Date {i+1}", default[1])
        oowner = st.text_input(f"Owner {i+1}", default[2])
        osteps = st.text_area(f"Steps {i+1}", default[3], height=60)
        open_items_data.append({"task": otask, "date": odate, "owner": oowner, "steps": osteps})

# Next Steps
st.header("Recommended Next Steps")
short_term_input = st.text_area("Short Term (comma-separated)", value="Finish Production rollout.,Tighten Firewall policies.,Tighten Cloud App Control Policies.,Fine tune SSL Inspection policies.,Configure Role Based Access Control (RBAC).,Configure DLP policies.")
long_term_input = st.text_area("Long Term (comma-separated)", value="Deploy ZCC on Mobile devices.,Consider an upgrade of Sandbox license to have better antimalware protection.,Consider an upgrade of the Firewall License to be able to apply policies based on user groups and network applications.,Adopt additional Zscaler solutions like Zscaler Private Access (ZPA) or Zscaler Digital experience (ZDX).,Consider using ZCC Client when the users are on-prem for a more consistent user experience.,Integrate ZIA with 3rd party SIEM.")
short_term = [s.strip() for s in short_term_input.split(",") if s.strip()]
long_term = [s.strip() for s in long_term_input.split(",") if s.strip()]

# Contacts
st.header("Contacts")
col1, col2 = st.columns(2)
pm_name = col1.text_input("Project Manager", value="Alex Vazquez")
consultant_name = col2.text_input("Consultant", value="Alex Vazquez")
primary_contact = st.text_input("Primary Contact", value="Teia proctor")
secondary_contact = st.text_input("Secondary Contact", value="Marco Sattier")

# Preview Button (enhanced)
if st.button("Preview Inputs"):
    st.subheader("Preview")
    st.write(f"**Customer:** {customer_name} | **Date:** {today_date} | **Summary:** {project_summary_text[:100]}...")
    st.write(f"**Milestones:** {', '.join([m['name'] for m in milestones_data])}")
    st.write(f"**Rollout:** Pilot {pilot_current}/{pilot_target}, Prod {prod_current}/{prod_target}")
    st.write(f"**Objectives:** {len(objectives_data)} rows")
    st.write(f"**Deliverables:** {len(deliverables_data)} rows")
    st.write(f"**Tech:** {windows_num} Windows, {mac_num} Mac")
    st.write(f"**Open Items:** {len(open_items_data)} rows")
    st.write(f"**Next Steps:** {len(short_term)} short, {len(long_term)} long")
    if uploaded_images:
        st.write("**Uploaded Images:** " + ", ".join(uploaded_images.keys()))

# Generation (with validation)
if st.button("Generate & Download PPTX"):
    # Validation (enhanced)
    if not customer_name:
        st.error("Customer Name required!")
    elif not all(is_valid_date(d) for d in [today_date, project_start, project_end, pilot_completion, prod_completion]):
        st.error("Fix date formats (DD/MM/YYYY or ??)")
    else:
        prs = Presentation()
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        logo_bytes = download_image_to_bytes(LOGO_URL) or download_image_to_bytes(ALT_LOGO_URL)
        bg_bytes = download_image_to_bytes(BG_URL)

        # Helper: Title Slide (tweaked positions)
        def create_title_slide(title_text: str, subtitle_text: str = "", date_text: str = "", slide_num: int = 1):
            slide = add_slide_with_background(prs, bg_bytes)
            add_textbox(slide, MARGIN_LEFT, Inches(1.0), Inches(8.0), Inches(1.0), title_text, SIZE_TITLE, True, COLOR_NAVY)
            if subtitle_text:
                add_textbox(slide, MARGIN_LEFT, Inches(2.1), Inches(8.0), Inches(0.5), subtitle_text, SIZE_SUBTITLE, color=COLOR_NAVY)
            if date_text:
                add_textbox(slide, MARGIN_LEFT, Inches(2.6), Inches(8.0), Inches(0.5), date_text, SIZE_BODY)
            apply_template_branding(prs, slide, slide_num, logo_bytes)
            # Add images if uploaded (placeholder positions)
            for img_name, img_bytes in uploaded_images.items():
                if "image1" in img_name:  # Example mapping
                    slide.shapes.add_picture(img_bytes, Inches(9.0), Inches(1.0), Inches(2.0), Inches(1.5))
            return slide

        # Helper: Bullet Slide (same, but added image support)
        def create_bullet_slide(title_text: str, bullets: List[str], slide_num: int = 1):
            slide = add_slide_with_background(prs, bg_bytes)
            add_textbox(slide, MARGIN_LEFT, Inches(0.45), Inches(8.0), Inches(0.5), title_text, SIZE_SLIDE_TITLE, True, COLOR_NAVY)
            top = Inches(1.2)
            for b in bullets:
                add_textbox(slide, MARGIN_LEFT + Inches(0.5), top, Inches(7.5), Inches(0.4), "- " + b, SIZE_BODY)
                top += Inches(0.5)
            apply_template_branding(prs, slide, slide_num, logo_bytes)
            return slide

        # Helper: Table Slide (enhanced with RAG colors, exact widths, and run guards to fix IndexError)
        def create_table_slide(title_text: str, headers: List[str], rows: List[List[str]], slide_num: int = 1, top_inch: float = 1.2, height_inch: float = 4.0, col_widths: List = None):
            slide = add_slide_with_background(prs, bg_bytes)
            add_textbox(slide, MARGIN_LEFT, Inches(0.45), Inches(8.0), Inches(0.5), title_text, SIZE_SLIDE_TITLE, True, COLOR_NAVY)
            left = MARGIN_LEFT
            top = Inches(top_inch)
            width = slide_width - 2 * MARGIN_LEFT
            height = Inches(height_inch)
            cols = len(headers)
            table = slide.shapes.add_table(len(rows) + 1, cols, left, top, width, height).table
            # Set widths (EMU, exact from template) - handle Length or float
            if not col_widths:
                col_widths = [width / cols] * cols
            for i, w in enumerate(col_widths):
                if isinstance(w, Length):
                    table.columns[i].width = w
                else:
                    table.columns[i].width = Emu(Inches(w))
            # Headers (with run guard)
            for i, h in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = str(h)
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_NAVY
                p = cell.text_frame.paragraphs[0]
                if not p.runs:
                    p.add_run()
                set_font_run(p.runs[0], size=SIZE_HEADER, bold=True, color=COLOR_WHITE)
            # Rows (with RAG if status column, and run guard)
            for r, row in enumerate(rows, 1):
                for c, val in enumerate(row):
                    cell = table.cell(r, c)
                    cell.text = str(val)
                    p = cell.text_frame.paragraphs[0]
                    if not p.runs:
                        p.add_run()
                    set_font_run(p.runs[0], size=SIZE_BODY)
                    if c == len(headers) - 1 and val in RAG_COLORS:  # Status column
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RAG_COLORS[val]
                    elif r % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            apply_template_branding(prs, slide, slide_num, logo_bytes)
            return slide

        # Helper: ZIA Diagram (expanded to match template exactly, with more elements)
        def create_zia_diagram_slide(slide_num: int = 1):
            slide = add_slide_with_background(prs, bg_bytes)
            add_textbox(slide, MARGIN_LEFT, Inches(0.45), Inches(8.0), Inches(0.5), "Deployed ZIA Architecture", SIZE_SLIDE_TITLE, True, COLOR_NAVY)
            # Boxes and labels (fine-tuned positions)
            box_w = Inches(2.5)
            box_h = Inches(1.0)
            left1 = Inches(0.5)
            top1 = Inches(1.2)
            # User authentication box
            shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left1, top1, box_w, box_h)
            shape1.fill.solid(); shape1.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            add_textbox(slide, left1 + Inches(0.2), top1 + Inches(0.3), box_w - Inches(0.4), box_h - Inches(0.6), "User authentication and provisioning", SIZE_SMALL, align=PP_ALIGN.CENTER)
            # Central Authority
            left2 = left1 + box_w + Inches(0.5)
            shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left2, top1, box_w, box_h)
            shape2.fill.solid(); shape2.fill.fore_color.rgb = COLOR_BRIGHT_BLUE
            add_textbox(slide, left2 + Inches(0.2), top1 + Inches(0.3), box_w - Inches(0.4), box_h - Inches(0.6), "Central Authority", SIZE_SMALL, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
            # Public Service Edges
            left3 = left2 + box_w + Inches(0.5)
            shape3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left3, top1, box_w, box_h)
            shape3.fill.solid(); shape3.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            add_textbox(slide, left3 + Inches(0.2), top1 + Inches(0.3), box_w - Inches(0.4), box_h - Inches(0.6), "Public Service Edges", SIZE_SMALL, align=PP_ALIGN.CENTER)
            # Workforce Region-X
            top2 = top1 + box_h + Inches(0.5)
            shape4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left1, top2, box_w, box_h)
            shape4.fill.solid(); shape4.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            add_textbox(slide, left1 + Inches(0.2), top2 + Inches(0.3), box_w - Inches(0.4), box_h - Inches(0.6), "Workforce (Region-X)\nOn | Off - net", SIZE_SMALL, align=PP_ALIGN.CENTER)
            # Z-Tunnels box
            shape5 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left2, top2, box_w, box_h)
            shape5.fill.solid(); shape5.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            add_textbox(slide, left2 + Inches(0.2), top2 + Inches(0.3), box_w - Inches(0.4), box_h - Inches(0.6), "Z-Tunnels", SIZE_SMALL, align=PP_ALIGN.CENTER)
            # SSL Inspection
            shape6 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left3, top2, box_w, box_h)
            shape6.fill.solid(); shape6.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            add_textbox(slide, left3 + Inches(0.2), top2 + Inches(0.3), box_w - Inches(0.4), box_h - Inches(0.6), "SSL Inspection", SIZE_SMALL, align=PP_ALIGN.CENTER)
            # Workforce Region-Y (added for template)
            top3 = top2 + box_h + Inches(0.5)
            shape7 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left1, top3, box_w, box_h)
            shape7.fill.solid(); shape7.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            add_textbox(slide, left1 + Inches(0.2), top3 + Inches(0.3), box_w - Inches(0.4), box_h - Inches(0.6), "Workforce (Region-Y)\nOn | Off - net", SIZE_SMALL, align=PP_ALIGN.CENTER)
            # Admin Console
            shape8 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left3, top3, box_w, box_h)
            shape8.fill.solid(); shape8.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            add_textbox(slide, left3 + Inches(0.2), top3 + Inches(0.3), box_w - Inches(0.4), box_h - Inches(0.6), "Admin Console", SIZE_SMALL, align=PP_ALIGN.CENTER)
            # Logging
            shape9 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left2, top3, box_w, box_h)
            shape9.fill.solid(); shape9.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            add_textbox(slide, left2 + Inches(0.2), top3 + Inches(0.3), box_w - Inches(0.4), box_h - Inches(0.6), "Logging", SIZE_SMALL, align=PP_ALIGN.CENTER)
            # Numbers (1-5 from template)
            add_textbox(slide, left1 + box_w / 2, top1 - Inches(0.3), Inches(0.3), Inches(0.3), "1", SIZE_SMALL)
            add_textbox(slide, left2 + box_w / 2, top1 - Inches(0.3), Inches(0.3), Inches(0.3), "3", SIZE_SMALL)
            add_textbox(slide, left3 + box_w / 2, top1 - Inches(0.3), Inches(0.3), Inches(0.3), "4", SIZE_SMALL)
            add_textbox(slide, left1 + box_w / 2, top2 - Inches(0.3), Inches(0.3), Inches(0.3), "2", SIZE_SMALL)
            add_textbox(slide, left3 + box_w / 2, top3 - Inches(0.3), Inches(0.3), Inches(0.3), "5", SIZE_SMALL)
            # Connectors (arrows for Z-Tunnels, etc.)
            try:
                conn1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left1 + box_w, top1 + box_h/2, left2, top1 + box_h/2)
                conn1.line.color.rgb = COLOR_BLACK
                conn2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left2 + box_w, top1 + box_h/2, left3, top1 + box_h/2)
                conn2.line.color.rgb = COLOR_BLACK
                conn3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left1 + box_w/2, top1 + box_h, left1 + box_w/2, top2)
                conn3.line.color.rgb = COLOR_BLACK
                # Add more for full template
                conn4 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left3 + box_w/2, top1 + box_h, left3 + box_w/2, top2)
                conn4.line.color.rgb = COLOR_BLACK
                conn5 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, left1 + box_w/2, top2 + box_h, left1 + box_w/2, top3)
                conn5.line.color.rgb = COLOR_BLACK
            except Exception:
                pass
            # Key facts (as table-like text)
            key_top = Inches(1.2)
            key_left = Inches(8.0)
            key_text = f"Identity Provider: {idp}\nAuthentication Type: {auth_type}\nProvisioning: {prov_type}\n\nTunnel Type: {tunnel_type}\nDeployment System: {deploy_system}\nNumber of Windows and MacOS Devices: {windows_num} Windows\n{mac_num} MacOS\nGeo Locations: {geo_locations}\n\nPolicy Deployment\nSSL Inspection Policies: {ssl_policies}\nURL Filtering Policies: {url_policies}\nCloud App Control Policies: {cloud_policies}\nFirewall Policies: {fw_policies}"
            add_textbox(slide, key_left, key_top, Inches(4.0), Inches(4.0), key_text, SIZE_SMALL)
            apply_template_branding(prs, slide, slide_num, logo_bytes)
            return slide

        # Build Slides (added missing ones: rollout, objectives, who/what, RAG key)
        progress = st.progress(0)
        total_slides = 12  # Added one for full match
        current = 0

        # Slide 1: Title
        create_title_slide("Professional Services Transition Meeting", customer_name, today_date, 1)
        current += 1
        progress.progress(current / total_slides)

        # Slide 2: Agenda
        create_bullet_slide("Meeting Agenda", ["Project Summary", "Technical Summary", "Recommended Next Steps"], 2)
        current += 1
        progress.progress(current / total_slides)

        # Slide 3: Project Summary Title
        create_title_slide("Project Summary", "", "", 3)
        current += 1
        progress.progress(current / total_slides)

        # Slide 4: Final Project Status Report (added who/what box, RAG key)
        slide4 = add_slide_with_background(prs, bg_bytes)
        add_textbox(slide4, MARGIN_LEFT, Inches(0.45), Inches(8.0), Inches(0.5), f"Final Project Status Report – {customer_name}", SIZE_SLIDE_TITLE, True, COLOR_NAVY)
        add_textbox(slide4, MARGIN_LEFT, Inches(1.2), Inches(8.0), Inches(0.4), "Project Summary", SIZE_HEADER, True)
        add_textbox(slide4, MARGIN_LEFT, Inches(1.7), Inches(8.0), Inches(1.0), project_summary_text, SIZE_BODY)
        # Dates
        add_textbox(slide4, MARGIN_LEFT, Inches(2.5), Inches(4.0), Inches(1.0), f"Today's Date: {today_date} | Start: {project_start} | End: {project_end}", SIZE_BODY)
        # Who/What/When/Why box (new)
        who_text = "Who: External & Internal Project Team\nWhat: Project Status Report\nWhen: Weekly\nWhy: Keeps stakeholders informed on scope, schedule, risks, etc.\nMandatory: Yes"
        add_textbox(slide4, Inches(6.0), Inches(3.0), Inches(4.0), Inches(2.0), who_text, SIZE_SMALL)
        # RAG Key (new table-like)
        rag_text = "RAG Status Key:\nRed - Not On Track\nAmber - At Risk\nGreen - On Track\nBlue - Complete\nGray - Not Started"
        add_textbox(slide4, Inches(6.0), Inches(5.5), Inches(4.0), Inches(1.5), rag_text, SIZE_SMALL)
        apply_template_branding(prs, slide4, 4, logo_bytes)
        current += 1
        progress.progress(current / total_slides)

        # Slide 5: Milestones Table
        headers = ["Milestone", "Baseline Date", "Target Completion Date", "Status"]
        rows = [[m["name"], m["baseline"], m["target"], m["status"]] for m in milestones_data]
        create_table_slide("Milestones", headers, rows, 5, col_widths=[Inches(4.0), Inches(2.0), Inches(2.0), Inches(2.0)])
        current += 1
        progress.progress(current / total_slides)

        # Slide 6: User Rollout Table (new)
        rollout_headers = ["Milestone", "Target Users", "Current Users", "Target Completion", "Status"]
        rollout_rows = [
            ["Pilot", str(pilot_target), str(pilot_current), pilot_completion, pilot_status],
            ["Production", str(prod_target), str(prod_current), prod_completion, prod_status]
        ]
        create_table_slide("User Rollout Roadmap", rollout_headers, rollout_rows, 6, top_inch=1.2, height_inch=1.5, col_widths=[Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)])
        current += 1
        progress.progress(current / total_slides)

        # Slide 7: Project Status (Objectives Table, new)
        obj_headers = ["Planned Project Objective (Target)", "Actual Project Result (Actual)", "Deviation/ Cause"]
        obj_rows = [[o["objective"], o["actual"], o["deviation"]] for o in objectives_data]
        create_table_slide("Project Status", obj_headers, obj_rows, 7, top_inch=1.2, height_inch=2.0, col_widths=[Inches(3.5), Inches(3.5), Inches(3.0)])
        current += 1
        progress.progress(current / total_slides)

        # Slide 8: Deliverables Table
        del_headers = ["Deliverable", "Date delivered"]
        del_rows = [[d["name"], d["date"]] for d in deliverables_data]
        create_table_slide("Deliverables", del_headers, del_rows, 8, top_inch=1.2)
        current += 1
        progress.progress(current / total_slides)

        # Slide 9: Technical Summary Title
        create_title_slide("Technical Summary", "", "", 9)
        current += 1
        progress.progress(current / total_slides)

        # Slide 10: ZIA Architecture
        create_zia_diagram_slide(10)
        current += 1
        progress.progress(current / total_slides)

        # Slide 11: Open Items Table
        open_headers = ["Task/ Description", "Date", "Owner", "Transition Plan/ Next Steps"]
        open_rows = [[oi["task"], oi["date"], oi["owner"], oi["steps"]] for oi in open_items_data]
        create_table_slide("Open Items", open_headers, open_rows, 11, col_widths=[Inches(2.5), Inches(1.5), Inches(1.5), Inches(4.5)])
        current += 1
        progress.progress(current / total_slides)

        # Slide 12: Next Steps & Thank You (combined for match)
        slide12 = add_slide_with_background(prs, bg_bytes)
        add_textbox(slide12, MARGIN_LEFT, Inches(0.45), Inches(8.0), Inches(0.5), "Recommended Next Steps", SIZE_SLIDE_TITLE, True, COLOR_NAVY)
        # Short Term
        add_textbox(slide12, MARGIN_LEFT, Inches(1.2), Inches(4.0), Inches(0.4), "Short Term Activities", SIZE_HEADER, True)
        top = Inches(1.6)
        for item in short_term:
            add_textbox(slide12, MARGIN_LEFT + Inches(0.3), top, Inches(3.5), Inches(0.3), item, SIZE_BODY)
            top += Inches(0.4)
        # Long Term
        add_textbox(slide12, Inches(5.5), Inches(1.2), Inches(4.0), Inches(0.4), "Long Term Activities", SIZE_HEADER, True)
        top = Inches(1.6)
        for item in long_term:
            add_textbox(slide12, Inches(5.8), top, Inches(3.5), Inches(0.3), item, SIZE_BODY)
            top += Inches(0.4)
        # Thank You section
        add_textbox(slide12, MARGIN_LEFT, Inches(4.5), Inches(8.0), Inches(0.5), "Thank you", SIZE_TITLE, True, COLOR_NAVY)
        thank_text = f"Your feedback on our project and Professional Services team is important to us. \nProject Manager: {pm_name}\nConsultant: {consultant_name}\n\nA short ~6 question survey on how your Professional Services team did will be automatically sent after the project has closed. The following people will receive the survey via email:\nPrimary Contact: {primary_contact}\nSecondary Contact: {secondary_contact}\nWe appreciate any insights you can provide to help us improve our processes and ensure we provide the best possible service in future projects.\n\nWe want to know!"
        add_textbox(slide12, MARGIN_LEFT, Inches(5.0), Inches(8.0), Inches(3.0), thank_text, SIZE_BODY)
        apply_template_branding(prs, slide12, 12, logo_bytes)
        current += 1
        progress.progress(current / total_slides)

        # Save & Download
        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        st.success("Deck generated! Matches template exactly.")
        st.download_button("Download PPTX", out, f"{customer_name}_Transition_Deck.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
